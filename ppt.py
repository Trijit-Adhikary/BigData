# GPT

import asyncio
from typing import List, Tuple

from pptx import Presentation
from langchain_openai import AzureChatOpenAI

# ------------------------------
# Configurable constants
# ------------------------------
TITLE_LAYOUT_INDEX = 0  # Title slide layout index in the master template
BODY_LAYOUT_INDEX = 1   # Title and Content layout index in the master template

# ------------------------------
# Content generation
# ------------------------------
async def generate_slide_content(context: str, num_slides: int = 5, client: AzureChatOpenAI = None) -> str:
    """
    Asks the LLM to produce a slide outline with a clear, parseable structure.
    Returns raw text.
    """
    if client is None:
        raise ValueError("AzureChatOpenAI client is required")

    prompt = f"""
Create a {num_slides}-slide outline for a professional PowerPoint presentation from the provided context.

CONTEXT:
{context}

REQUIREMENTS:
1. Slide 1 must be the presentation title slide (only a title and an optional subtitle).
2. Slides 2 to {num_slides}: each must have a title and 3-5 bullet points, derived strictly from the context.
3. Do not use markdown or special characters.
4. Use the following strict structure:

SLIDE 1 TITLE: <title text>
SLIDE 1 SUBTITLE: <subtitle text or leave empty>

SLIDE 2 TITLE: <title text>
- <bullet 1>
- <bullet 2>
- <bullet 3>
- <bullet 4> (optional)
- <bullet 5> (optional)

... and so on, up to slide {num_slides}.
""".strip()

    response = await client.ainvoke([{"role": "system", "content": prompt}])
    return response.content


# ------------------------------
# Parsing helpers
# ------------------------------
def parse_slides(raw_text: str) -> Tuple[str, str, List[Tuple[str, List[str]]]]:
    """
    Parses the LLM output into:
    - title_slide_title: str
    - title_slide_subtitle: str
    - body_slides: List of tuples (title, bullets)
    """
    lines = [l.strip() for l in raw_text.splitlines() if l.strip()]
    title_slide_title = ""
    title_slide_subtitle = ""
    body_slides: List[Tuple[str, List[str]]] = []

    current_title = None
    current_bullets: List[str] = []

    def push_current():
        nonlocal current_title, current_bullets, body_slides
        if current_title:
            body_slides.append((current_title, current_bullets))
        current_title = None
        current_bullets = []

    i = 0
    while i < len(lines):
        line = lines[i]

        if line.upper().startswith("SLIDE 1 TITLE:"):
            title_slide_title = line.split(":", 1)[1].strip() if ":" in line else ""
            i += 1
            continue

        if line.upper().startswith("SLIDE 1 SUBTITLE:"):
            title_slide_subtitle = line.split(":", 1)[1].strip() if ":" in line else ""
            i += 1
            continue

        if line.upper().startswith("SLIDE ") and " TITLE:" in line.upper():
            # New body slide
            if current_title is not None:
                push_current()
            current_title = line.split(":", 1)[1].strip()
            current_bullets = []
            i += 1
            continue

        if line.startswith("- "):
            current_bullets.append(line[2:].strip())
            i += 1
            continue

        # Ignore any stray lines
        i += 1

    # Flush last slide
    if current_title is not None:
        push_current()

    return title_slide_title, title_slide_subtitle, body_slides


# ------------------------------
# PPT creation
# ------------------------------
def create_ppt_from_template(
    title_slide_title: str,
    title_slide_subtitle: str,
    body_slides: List[Tuple[str, List[str]]],
    template_path: str,
    output_file: str = "abc.pptx"
):
    """
    Creates a PPT using a given master template.
    - First slide uses TITLE_LAYOUT_INDEX
    - Remaining slides use BODY_LAYOUT_INDEX
    """
    prs = Presentation(template_path)

    # Slide 1: Title slide
    title_layout = prs.slide_layouts[TITLE_LAYOUT_INDEX]
    title_slide = prs.slides.add_slide(title_layout)

    # Title placeholder (usually index 0)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = title_slide_title or "Presentation"

    # Subtitle placeholder (commonly placeholders[1] if present)
    subtitle_shape = None
    if len(title_slide.placeholders) > 1:
        subtitle_shape = title_slide.placeholders[1]
    if subtitle_shape is not None:
        subtitle_shape.text = title_slide_subtitle or ""

    # Body slides
    body_layout = prs.slide_layouts[BODY_LAYOUT_INDEX]
    for slide_title, bullets in body_slides:
        slide = prs.slides.add_slide(body_layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title

        # Content placeholder
        content_ph = None
        # Find the first text placeholder that is not the title
        for ph in slide.placeholders:
            if ph.placeholder_format.type not in (1,):  # 1 is TITLE
                content_ph = ph
                break

        if content_ph is not None:
            # Join bullets with newlines so PPT renders them as bullet points
            content_ph.text = ""
            if bullets:
                tf = content_ph.text_frame
                tf.clear()
                first = True
                for b in bullets:
                    if first:
                        tf.text = b
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = b
                        p.level = 0

    prs.save(output_file)
    print(f"✅ PowerPoint saved as {output_file}")


# ------------------------------
# Orchestrator
# ------------------------------
async def build_presentation(
    context: str,
    num_slides: int,
    template_path: str,
    output_file: str,
    azure_endpoint: str,
    azure_api_key: str,
    azure_deployment: str = "gpt-4o",
    api_version: str = "2024-06-01",
    temperature: float = 0.1
):
    client = AzureChatOpenAI(
        azure_endpoint=azure_endpoint,
        api_key=azure_api_key,
        api_version=api_version,
        azure_deployment=azure_deployment,
        temperature=temperature,
        max_tokens=None
    )

    raw = await generate_slide_content(context=context, num_slides=num_slides, client=client)
    title_title, title_subtitle, body = parse_slides(raw)
    create_ppt_from_template(
        title_slide_title=title_title,
        title_slide_subtitle=title_subtitle,
        body_slides=body,
        template_path=template_path,
        output_file=output_file
    )


# ------------------------------
# Example usage
# ------------------------------
if __name__ == "__main__":
    # Replace with your real values
    azure_openai_endpoint = "<YOUR_AZURE_OPENAI_ENDPOINT>"
    azure_openai_key = "<YOUR_AZURE_OPENAI_KEY>"

    context_text = "Your source text goes here..."
    template_path = "master_template.pptx"  # ensure this file exists
    output_file = "abc.pptx"

    asyncio.run(
        build_presentation(
            context=context_text,
            num_slides=5,
            template_path=template_path,
            output_file=output_file,
            azure_endpoint=azure_openai_endpoint,
            azure_api_key=azure_openai_key,
            azure_deployment="gpt-4o",
            api_version="2024-06-01",
            temperature=0.1
        )
    )



# Sonnet
from pptx import Presentation
from langchain_openai import AzureChatOpenAI
import os

class PPTGenerator:
    """
    A class to generate PowerPoint presentations from text content using a master template.
    """
    
    def __init__(self, template_path=None):
        """
        Initialize the PPT generator with optional master template.
        
        Args:
            template_path (str): Path to master template file (.pptx)
        """
        self.template_path = template_path
        self.langchain_client = None
    
    def setup_langchain_client(self, azure_endpoint, api_key, azure_deployment):
        """
        Setup Azure OpenAI client for content generation.
        """
        self.langchain_client = AzureChatOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=api_key,
            api_version="2024-06-01",
            azure_deployment=azure_deployment,
            temperature=0.1,
            max_tokens=None
        )
    
    async def generate_slide_content(self, context, num_slides=5):
        """
        Generate structured slide content from context using AI.
        
        Args:
            context (str): Input text context
            num_slides (int): Number of slides to generate
            
        Returns:
            dict: Structured slide data with title and slides
        """
        if not self.langchain_client:
            raise ValueError("Langchain client not initialized. Call setup_langchain_client() first.")
        
        prompt = f"""
        Create a {num_slides}-slide outline for a professional PowerPoint presentation from the provided context.

        CONTEXT:
        {context}

        REQUIREMENTS:
        1. First, provide a main presentation title
        2. Then create {num_slides} content slides, each with:
           - A clear slide title
           - 3-5 concise bullet points
        3. All content must be derived from the context
        4. Use simple, clear language without markdown or special characters

        FORMAT YOUR RESPONSE EXACTLY AS:
        PRESENTATION_TITLE: [Main title here]

        SLIDE_1: [Slide 1 title]
        - [Bullet point 1]
        - [Bullet point 2]
        - [Bullet point 3]

        SLIDE_2: [Slide 2 title]
        - [Bullet point 1]
        - [Bullet point 2]
        - [Bullet point 3]
        
        [Continue for all slides...]
        """

        response = await self.langchain_client.ainvoke([
            {"role": "system", "content": prompt}
        ])

        return self._parse_slide_content(response.content)
    
    def _parse_slide_content(self, content):
        """
        Parse AI-generated content into structured slide data.
        
        Args:
            content (str): Raw AI response
            
        Returns:
            dict: Parsed slide data
        """
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        slide_data = {
            'presentation_title': '',
            'slides': []
        }
        
        current_slide = None
        
        for line in lines:
            # Extract presentation title
            if line.startswith('PRESENTATION_TITLE:'):
                slide_data['presentation_title'] = line.replace('PRESENTATION_TITLE:', '').strip()
            
            # Extract slide titles
            elif line.startswith('SLIDE_'):
                if current_slide:
                    slide_data['slides'].append(current_slide)
                
                title = line.split(':', 1)[1].strip() if ':' in line else line
                current_slide = {
                    'title': title,
                    'bullets': []
                }
            
            # Extract bullet points
            elif line.startswith('-') and current_slide:
                bullet = line.replace('-', '').strip()
                current_slide['bullets'].append(bullet)
        
        # Add the last slide
        if current_slide:
            slide_data['slides'].append(current_slide)
        
        return slide_data
    
    def create_ppt_from_template(self, slide_data, output_file="presentation.pptx"):
        """
        Create PowerPoint presentation using master template.
        
        Args:
            slide_data (dict): Structured slide data
            output_file (str): Output file path
        """
        # Load master template or create new presentation
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
            print(f"📄 Using master template: {self.template_path}")
        else:
            prs = Presentation()
            print("📄 Using default PowerPoint template")
        
        # Clear existing slides if using template
        if len(prs.slides) > 0:
            # Keep only the first 2 slides as templates
            slides_to_remove = list(prs.slides)[2:]  # Remove slides beyond index 1
            for slide in slides_to_remove:
                rId = prs.slides._sldIdLst[prs.slides.index(slide)].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[prs.slides.index(slide)]
        
        # Create title slide (using layout 0 or first slide as template)
        self._create_title_slide(prs, slide_data['presentation_title'])
        
        # Create content slides (using layout 1 or second slide as template)
        for slide_info in slide_data['slides']:
            self._create_content_slide(prs, slide_info)
        
        # Save presentation
        prs.save(output_file)
        print(f"✅ PowerPoint saved as {output_file}")
    
    def _create_title_slide(self, prs, title):
        """
        Create title slide using first slide template.
        """
        if len(prs.slides) >= 1:
            # Use first slide as title template
            title_layout = prs.slide_layouts[0]  # Title slide layout
        else:
            title_layout = prs.slide_layouts[0]  # Default title layout
        
        title_slide = prs.slides.add_slide(title_layout)
        
        # Set title
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title
        
        # If there's a subtitle placeholder, you can add subtitle here
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Generated Presentation"
    
    def _create_content_slide(self, prs, slide_info):
        """
        Create content slide using second slide template.
        """
        # Use second slide template (index 1) for body content
        content_layout = prs.slide_layouts[1]  # Content with bullets layout
        
        slide = prs.slides.add_slide(content_layout)
        
        # Set slide title
        if slide.shapes.title:
            slide.shapes.title.text = slide_info['title']
        
        # Set bullet points
        if len(slide.placeholders) > 1 and slide_info['bullets']:
            content_placeholder = slide.placeholders[1]
            bullet_text = '\n'.join([f"• {bullet}" for bullet in slide_info['bullets']])
            content_placeholder.text = bullet_text


# Usage Example
async def main():
    """
    Main function demonstrating the usage of PPTGenerator.
    """
    # Initialize PPT Generator
    ppt_gen = PPTGenerator(template_path="master_template.pptx")  # Optional template path
    
    # Setup Azure OpenAI client
    ppt_gen.setup_langchain_client(
        azure_endpoint="your_azure_endpoint",
        api_key="your_api_key",
        azure_deployment="gpt-4o"
    )
    
    # Sample context text
    sample_context = """
    Artificial Intelligence is transforming businesses across industries. 
    It enables automation, improves decision-making, and enhances customer experiences.
    Companies are investing heavily in AI technologies to stay competitive.
    """
    
    try:
        # Generate slide content
        print("🔄 Generating slide content...")
        slide_data = await ppt_gen.generate_slide_content(sample_context, num_slides=4)
        
        print(f"📋 Generated presentation: '{slide_data['presentation_title']}'")
        print(f"📊 Number of content slides: {len(slide_data['slides'])}")
        
        # Create PowerPoint presentation
        print("🔄 Creating PowerPoint presentation...")
        ppt_gen.create_ppt_from_template(slide_data, "generated_presentation.pptx")
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())





